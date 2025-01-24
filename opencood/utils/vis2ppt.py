from pptx import Presentation
from pptx.util import Inches
import os
from PIL import Image
import numpy as np
from tqdm import tqdm

cache = "/home/xiangbog/Folder/AdapterCAV/opencood/utils/cache.png"


def crop_detect(detect):
    with Image.open(detect) as img:
        img_np = np.array(img)
        H, W = img_np.shape[:2]
        H_new = H / 2 * (1 - 0.4 / 7.06)
        W_new = W / 2 * (1 - 0.4 / 7.06)
        boundH = (H - H_new) / 2
        boundW = (W - W_new) / 2
        img_np = img_np[int(boundH) : int(boundH + H_new), int(boundW) : int(boundW + W_new)]
        img = Image.fromarray(img_np)
        img.save(cache)


def crop_feat(feat):
    with Image.open(feat) as img:
        img_np = np.array(img)
        H, W = img_np.shape[:2]
        H_new = H / 2
        W_new = W / 2
        boundH = (H - H_new) / 2
        boundW = (W - W_new) / 2
        img_np = img_np[int(boundH) : int(boundH + H_new), int(boundW) : int(boundW + W_new)]
        img = Image.fromarray(img_np)
        img.save(cache)


# Create a presentation object
prs = Presentation()

try:
    for sid in tqdm(list(range(0, 1000, 40))):

        # Add a slide layout (0 for title slide, 1 for title and content)
        slide_layout = prs.slide_layouts[5]  # Using layout 5 for blank slide
        slide = prs.slides.add_slide(slide_layout)


        # Define the image paths and positions
        root = "/home/xiangbog/Folder/AdapterCAV/opencood/logs/OURS_4Agents_ObjDetect/final_infer"

        adapter = "convnext_crop_w_output"
        identity = "identity_m2"

        dir = "vis_intermediate_102.4_102.4_epoch1"

        dir_ego = os.path.join(root, adapter, dir)
        dir_identity = os.path.join(root, identity, dir)
        idx = f"{sid:05}"

        images = []
        image_size = 2
        image_spacing = 0.1
        image_spacing_large = 0.5
        initial_x = 0.5
        initial_y = 1
        curr_x = initial_x
        curr_y = initial_y

        # # Iterate over the images with prefix "bev_bev_{idx}M" in the folder dir_ego
        M_prefix = f"bev_{idx}M"
        M2P_prefix = f"bev_{idx}M2P"
        M2P2M_prefix = f"bev_{idx}M2P2M"
        for img in os.listdir(dir_ego):
            if img.startswith(M_prefix + "_"):
                img = img[len(M_prefix) :]
                Midx = f"{int(img.split('.')[0][-1])}"
                M = os.path.join(dir_ego, M_prefix + img)
                M2P = os.path.join(dir_ego, M2P_prefix + "_" + Midx + ".png")
                M2P2M = os.path.join(dir_ego, M2P2M_prefix + "_" + Midx + ".png")

                images.append((M, crop_feat, curr_x, curr_y, image_size, image_size))
                images.append((M2P, crop_feat, curr_x, curr_y + image_size + image_spacing_large + image_spacing, image_size, image_size))
                images.append(
                    (M2P2M, crop_feat, curr_x, curr_y + 2 * (image_size + image_spacing_large + image_spacing), image_size, image_size)
                )
                curr_x = curr_x + image_size + image_spacing

        curr_x = curr_x + image_spacing_large

        FO = os.path.join(dir_identity, f"bev_{idx}FM.png")
        images.append((FO, crop_feat, curr_x, initial_y, image_size, image_size))

        FM2P = os.path.join(dir_ego, f"bev_{idx}FM2P_.png")
        images.append((FM2P, crop_feat, curr_x, initial_y + image_size + image_spacing_large, image_size, image_size))

        FM = os.path.join(dir_ego, f"bev_{idx}FM.png")
        images.append((FM, crop_feat, curr_x, initial_y + 2 * (image_size + image_spacing_large), image_size, image_size))

        curr_x = curr_x + image_spacing_large + image_size + image_spacing

        detect_original = os.path.join(dir_identity, f"bev_{idx}.png")
        images.append((detect_original, crop_detect, curr_x, initial_y, image_size, image_size))

        detect_protocol = os.path.join(dir_ego + "_protocol", f"bev_{idx}.png")
        images.append(
            (detect_protocol, crop_detect, curr_x, initial_y + image_size + image_spacing_large, image_size, image_size)
        )

        detect = os.path.join(dir_ego, f"bev_{idx}.png")
        images.append((detect, crop_detect, curr_x, initial_y + 2 * (image_size + image_spacing_large), image_size, image_size))


        # Add images to the slide
        for img_path, crop_func, left, top, width, height in images:
            crop_func(img_path)

            slide.shapes.add_picture(cache, Inches(left), Inches(top), Inches(width), Inches(height))

        # # Save the presentation
        # if not os.path.exists(os.path.join(root, adapter, "generated_presentation")):
        #     os.mkdir(os.path.join(root, adapter, "generated_presentation"))
except:
    pass
prs.save(os.path.join(root, adapter, f"generated_presentation.pptx"))
